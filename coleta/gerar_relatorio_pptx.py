"""
Gerador de relatorio executivo em PowerPoint.
==============================================
Gera um PPT por regional com resumo, breves lancamentos e destaques.

Uso:
    python gerar_relatorio_pptx.py
    python gerar_relatorio_pptx.py --regional SP_SPRM
    python gerar_relatorio_pptx.py --saida relatorio_custom.pptx
"""

import os
import sys
import argparse
from datetime import datetime
from collections import defaultdict, Counter

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from config.settings import DATA_DIR
from config.regionais import classificar_regional, nome_regional, REGIONAIS
from data.database import get_connection, obter_changelog
from detectar_destaques import detectar_destaques, agrupar_por_regional, carregar_empreendimentos


# ============================================================
# CONSTANTES DE ESTILO
# ============================================================

COR_TITULO = RGBColor(0x1A, 0x52, 0x76)   # Azul escuro
COR_SUBTITULO = RGBColor(0x2C, 0x3E, 0x50)
COR_HEADER = RGBColor(0x1A, 0x52, 0x76)
COR_HEADER_TEXTO = RGBColor(0xFF, 0xFF, 0xFF)
COR_TEXTO = RGBColor(0x33, 0x33, 0x33)
COR_CINZA = RGBColor(0x99, 0x99, 0x99)

TEMPLATE_PATH = os.path.join(os.path.dirname(os.path.abspath(__file__)), "config", "template_relatorio.pptx")

# Slide dimensions A4 retrato (Emu)
SLIDE_WIDTH = Emu(6858000)    # 7.5" aprox (A4 width)
SLIDE_HEIGHT = Emu(9144000)   # 10" aprox (A4 height)


def criar_apresentacao():
    """Cria apresentacao, usando template se existir."""
    if os.path.exists(TEMPLATE_PATH):
        prs = Presentation(TEMPLATE_PATH)
    else:
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT
    return prs


def adicionar_texto(slide, left, top, width, height, texto, tamanho=12,
                    negrito=False, cor=COR_TEXTO, alinhamento=PP_ALIGN.LEFT):
    """Adiciona caixa de texto ao slide."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = texto
    p.font.size = Pt(tamanho)
    p.font.bold = negrito
    p.font.color.rgb = cor
    p.alignment = alinhamento
    return tf


def adicionar_tabela(slide, left, top, width, headers, rows, col_widths=None):
    """Adiciona tabela ao slide."""
    n_rows = len(rows) + 1  # +1 header
    n_cols = len(headers)

    table_shape = slide.shapes.add_table(n_rows, n_cols,
                                          Inches(left), Inches(top),
                                          Inches(width), Inches(0.3 * n_rows))
    table = table_shape.table

    # Larguras de coluna
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    # Header
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(9)
            paragraph.font.bold = True
            paragraph.font.color.rgb = COR_HEADER_TEXTO
        cell.fill.solid()
        cell.fill.fore_color.rgb = COR_HEADER

    # Dados
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val) if val is not None else ""
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(8)
                paragraph.font.color.rgb = COR_TEXTO

    return table


def slide_capa(prs, regional_nome):
    """Adiciona slide de capa."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout em branco

    adicionar_texto(slide, 0.5, 2.5, 6.5, 1,
                    "Inteligencia de Mercado", tamanho=28,
                    negrito=True, cor=COR_TITULO, alinhamento=PP_ALIGN.CENTER)

    adicionar_texto(slide, 0.5, 3.5, 6.5, 0.8,
                    f"Regional: {regional_nome}", tamanho=20,
                    cor=COR_SUBTITULO, alinhamento=PP_ALIGN.CENTER)

    data = datetime.now().strftime("%B %Y").title()
    adicionar_texto(slide, 0.5, 4.5, 6.5, 0.5,
                    data, tamanho=14,
                    cor=COR_CINZA, alinhamento=PP_ALIGN.CENTER)


def slide_resumo(prs, regional_nome, empreendimentos):
    """Adiciona slide com resumo da regional."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    adicionar_texto(slide, 0.3, 0.2, 6.5, 0.5,
                    f"Resumo - {regional_nome}", tamanho=18,
                    negrito=True, cor=COR_TITULO)

    # Contagens por fase
    fases = Counter(e.get("fase", "Sem info") for e in empreendimentos)
    fases_texto = "\n".join(f"  {fase}: {count}" for fase, count in fases.most_common())

    # Contagens por empresa
    empresas = Counter(e.get("empresa", "?") for e in empreendimentos)
    empresas_texto = "\n".join(f"  {emp}: {count}" for emp, count in empresas.most_common(10))

    total = len(empreendimentos)
    com_preco = sum(1 for e in empreendimentos if e.get("preco_a_partir") and e["preco_a_partir"] > 0)
    com_coords = sum(1 for e in empreendimentos if e.get("latitude") and e.get("longitude"))

    resumo = (
        f"Total de empreendimentos: {total}\n"
        f"Com preco: {com_preco} ({100*com_preco//max(total,1)}%)\n"
        f"Geocodificados: {com_coords} ({100*com_coords//max(total,1)}%)\n"
        f"\nPor fase:\n{fases_texto}\n"
        f"\nPor empresa (top 10):\n{empresas_texto}"
    )

    adicionar_texto(slide, 0.3, 0.9, 6.5, 8, resumo, tamanho=10)


def slide_breves_lancamentos(prs, regional_nome, empreendimentos):
    """Adiciona slide(s) com breves lancamentos."""
    fases_breve = {"Breve Lançamento", "Futuro Lançamento", "Breve lançamento", "Futuro lançamento"}
    breves = [e for e in empreendimentos if e.get("fase") in fases_breve]

    if not breves:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    adicionar_texto(slide, 0.3, 0.2, 6.5, 0.5,
                    f"Breves Lancamentos - {regional_nome}", tamanho=18,
                    negrito=True, cor=COR_TITULO)

    headers = ["Empresa", "Nome", "Cidade", "Tipologia", "Preco"]
    rows = []
    for e in breves[:15]:  # Limitar a 15 por slide
        preco = f"R$ {e['preco_a_partir']:,.0f}" if e.get("preco_a_partir") else "-"
        rows.append([
            e.get("empresa", ""),
            e.get("nome", ""),
            e.get("cidade", ""),
            e.get("dormitorios_descricao", "-"),
            preco,
        ])

    adicionar_tabela(slide, 0.2, 0.9, 7, headers, rows,
                     col_widths=[1.2, 1.8, 1.2, 1.4, 1.0])


def slide_destaques(prs, regional_nome, destaques):
    """Adiciona slide(s) com destaques detectados."""
    if not destaques:
        return

    # Agrupar por tipo
    por_tipo = defaultdict(list)
    for d in destaques:
        por_tipo[d["tipo"]].append(d)

    for tipo, items in por_tipo.items():
        if not items:
            continue

        slide = prs.slides.add_slide(prs.slide_layouts[6])
        titulo_tipo = tipo.replace("_", " ").title()
        adicionar_texto(slide, 0.3, 0.2, 6.5, 0.5,
                        f"Destaques: {titulo_tipo} - {regional_nome}", tamanho=16,
                        negrito=True, cor=COR_TITULO)

        headers = ["Empresa", "Nome", "Cidade", "Detalhe"]
        rows = []
        for item in items[:12]:  # Limitar por slide
            rows.append([
                item.get("empresa", ""),
                item.get("nome", ""),
                item.get("cidade", ""),
                item.get("detalhe", ""),
            ])

        adicionar_tabela(slide, 0.2, 0.9, 7, headers, rows,
                         col_widths=[1.0, 1.5, 1.0, 3.0])


def gerar_relatorio(regional_filtro=None, saida=None):
    """Gera relatorio PPT completo."""
    prs = criar_apresentacao()

    # Carregar dados
    empreendimentos = carregar_empreendimentos()
    grupos = agrupar_por_regional(empreendimentos)
    destaques_todos = detectar_destaques(regional_filtro)

    regionais_ordenadas = sorted(grupos.keys(), key=lambda r: len(grupos[r]), reverse=True)

    for regional in regionais_ordenadas:
        if regional_filtro and regional != regional_filtro:
            continue

        emps = grupos[regional]
        regional_nome = nome_regional(regional)
        destaques = destaques_todos.get(regional, {}).get("destaques", [])

        # Slides da regional
        slide_capa(prs, regional_nome)
        slide_resumo(prs, regional_nome, emps)
        slide_breves_lancamentos(prs, regional_nome, emps)
        slide_destaques(prs, regional_nome, destaques)

    # Salvar
    if not saida:
        data = datetime.now().strftime("%Y%m%d")
        sufixo = f"_{regional_filtro}" if regional_filtro else ""
        saida = os.path.join(DATA_DIR, f"relatorio_im_{data}{sufixo}.pptx")

    prs.save(saida)
    print(f"Relatorio salvo: {saida}")
    return saida


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Gera relatorio executivo em PPT")
    parser.add_argument("--regional", help="Filtrar por regional (ex: SP_SPRM)")
    parser.add_argument("--saida", help="Caminho do arquivo de saida")
    args = parser.parse_args()

    gerar_relatorio(regional_filtro=args.regional, saida=args.saida)
